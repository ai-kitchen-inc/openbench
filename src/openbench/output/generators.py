"""Output format generators.

Provides concrete implementations of OutputGenerator for various formats:
- PDFGenerator: Generate PDF reports using ReportLab
- PowerPointGenerator: Generate PPTX presentations
- DashboardGenerator: Generate interactive dashboards
- AudioGenerator: Generate audio content from text
- MarkdownGenerator: Generate markdown files
"""

import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from openbench.core.abstractions import GeneratedOutput, OutputGenerator

logger = logging.getLogger(__name__)


class PDFGenerator(OutputGenerator):
    """
    Generate PDF reports using ReportLab.

    Implements the OutputGenerator interface for PDF output.
    Supports multiple templates and content types.

    Example:
        ```python
        from openbench.output.generators import PDFGenerator

        generator = PDFGenerator(template="report")
        result = generator.generate(
            content="This is my report content...",
            output_path="report.pdf",
            title="My Report"
        )
        print(result.file_path)  # report.pdf
        ```

    Content Types Supported:
        - str: Rendered as paragraphs
        - dict: Rendered as structured sections (keys as headings)
        - list: Rendered as bullet points
        - Dict from IntelligenceLayer: Extracts 'content' key
    """

    # Page sizes mapping
    PAGE_SIZES = {
        "letter": (612, 792),  # 8.5 x 11 inches
        "a4": (595.27, 841.89),  # 210 x 297 mm
        "legal": (612, 1008),  # 8.5 x 14 inches
    }

    def __init__(
        self,
        template: str = "default",
        page_size: str = "letter",
        margins: Optional[Dict[str, float]] = None,
        font_name: str = "Helvetica",
        font_size: int = 11,
        title_font_size: int = 18,
        heading_font_size: int = 14,
    ):
        """
        Initialize PDF generator.

        Args:
            template: Template name ('default', 'report', 'corporate')
            page_size: Page size ('letter', 'a4', 'legal')
            margins: Page margins in points {'top': 72, 'bottom': 72, 'left': 72, 'right': 72}
            font_name: Base font name
            font_size: Body text font size
            title_font_size: Title font size
            heading_font_size: Section heading font size
        """
        self.template = template
        self.page_size = page_size
        self.margins = margins or {"top": 72, "bottom": 72, "left": 72, "right": 72}
        self.font_name = font_name
        self.font_size = font_size
        self.title_font_size = title_font_size
        self.heading_font_size = heading_font_size

        logger.debug(f"PDFGenerator initialized (template: {template}, page_size: {page_size})")

    @property
    def output_format(self) -> str:
        """Output format identifier."""
        return "pdf"

    def validate(self, content: Any) -> bool:
        """
        Validate that content can be rendered as PDF.

        Args:
            content: Content to validate

        Returns:
            True if content can be rendered
        """
        if content is None:
            return False
        # Accept strings, dicts, lists, or objects with __str__
        return isinstance(content, (str, dict, list)) or hasattr(content, "__str__")

    def _extract_content(self, content: Any) -> str:
        """
        Extract text content from various input formats.

        Args:
            content: Input content (dict, str, list, etc.)

        Returns:
            Extracted text content
        """
        if isinstance(content, str):
            return content

        if isinstance(content, list):
            return "\n".join(f"- {item}" for item in content)

        if hasattr(content, "content"):
            return str(content.content)

        if hasattr(content, "output"):
            return str(content.output)

        if not isinstance(content, dict):
            return str(content)

        # Handle dict inputs from various layers
        if "content" in content:
            return str(content["content"])

        if "intelligence_output" in content:
            return self._extract_from_intelligence_output(content["intelligence_output"])

        if "raw_data" in content:
            return self._extract_from_raw_data(content["raw_data"])

        # Convert remaining dict keys to formatted text
        parts = [
            f"{key}:\n{value}"
            for key, value in content.items()
            if key not in ("metadata", "tokens_used", "model")
        ]
        return "\n\n".join(parts) if parts else str(content)

    def _extract_from_intelligence_output(self, output: Any) -> str:
        """Extract content from intelligence_output field."""
        if isinstance(output, dict) and "content" in output:
            return str(output["content"])

        if hasattr(output, "output"):
            return str(output.output)

        return str(output)

    def _extract_from_raw_data(self, raw: Any) -> str:
        """Extract content from raw_data field."""
        if isinstance(raw, list):
            return "\n\n".join(
                str(item.content) if hasattr(item, "content") else str(item)
                for item in raw
            )

        if hasattr(raw, "content"):
            return str(raw.content)

        return str(raw)

    def generate(
        self,
        content: Any,
        template: Optional[str] = None,
        output_path: str = "report.pdf",
        title: Optional[str] = None,
        author: Optional[str] = None,
        **options,
    ) -> GeneratedOutput:
        """
        Generate PDF report.

        Args:
            content: Content to render into PDF
            template: Template override (uses instance template if None)
            output_path: Output file path
            title: Document title (optional)
            author: Document author (optional)
            **options: Additional PDF-specific options

        Returns:
            GeneratedOutput with file path and metadata
        """
        used_template = template or self.template
        logger.info(f"Generating PDF: {output_path}")

        # Extract text content
        text_content = self._extract_content(content)

        # Ensure output directory exists
        output_dir = Path(output_path).parent
        if output_dir and str(output_dir) != ".":
            output_dir.mkdir(parents=True, exist_ok=True)

        # Try to use ReportLab
        try:
            size_bytes = self._generate_with_reportlab(
                text_content, output_path, title, author, used_template, **options
            )
        except ImportError:
            logger.warning("ReportLab not installed, falling back to simple text PDF")
            size_bytes = self._generate_simple_pdf(text_content, output_path)

        logger.info(f"PDF generated: {output_path} ({size_bytes} bytes)")

        return GeneratedOutput(
            file_path=output_path,
            format=self.output_format,
            size_bytes=size_bytes,
            metadata={
                "template": used_template,
                "page_size": self.page_size,
                "title": title,
                "author": author,
                "content_length": len(text_content),
                **options,
            },
        )

    def _generate_with_reportlab(
        self,
        content: str,
        output_path: str,
        title: Optional[str],
        author: Optional[str],
        template: str,
        **options,
    ) -> int:
        """
        Generate PDF using ReportLab.

        Args:
            content: Text content to render
            output_path: Output file path
            title: Document title
            author: Document author
            template: Template name
            **options: Additional options

        Returns:
            File size in bytes
        """
        from reportlab.lib.pagesizes import letter, A4, legal
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            SimpleDocTemplate,
            Paragraph,
            Spacer,
            PageBreak,
        )
        from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY

        # Select page size
        page_size_map = {"letter": letter, "a4": A4, "legal": legal}
        page_size = page_size_map.get(self.page_size.lower(), letter)

        # Create document
        doc = SimpleDocTemplate(
            output_path,
            pagesize=page_size,
            topMargin=self.margins["top"],
            bottomMargin=self.margins["bottom"],
            leftMargin=self.margins["left"],
            rightMargin=self.margins["right"],
            title=title or "OpenBench Report",
            author=author or "OpenBench",
        )

        # Get styles
        styles = getSampleStyleSheet()

        # Create custom styles
        title_style = ParagraphStyle(
            "CustomTitle",
            parent=styles["Heading1"],
            fontSize=self.title_font_size,
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName=f"{self.font_name}-Bold" if self.font_name == "Helvetica" else self.font_name,
        )

        heading_style = ParagraphStyle(
            "CustomHeading",
            parent=styles["Heading2"],
            fontSize=self.heading_font_size,
            spaceBefore=20,
            spaceAfter=10,
            fontName=f"{self.font_name}-Bold" if self.font_name == "Helvetica" else self.font_name,
        )

        body_style = ParagraphStyle(
            "CustomBody",
            parent=styles["Normal"],
            fontSize=self.font_size,
            leading=self.font_size * 1.4,
            alignment=TA_JUSTIFY,
            fontName=self.font_name,
        )

        # Build story (document content)
        story = []

        # Add title if provided
        if title:
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 20))

        # Add date and author for report template
        if template in ("report", "corporate"):
            date_str = datetime.now().strftime("%B %d, %Y")
            meta_style = ParagraphStyle(
                "Meta",
                parent=styles["Normal"],
                fontSize=10,
                textColor="gray",
                alignment=TA_CENTER,
            )
            if author:
                story.append(Paragraph(f"Author: {author}", meta_style))
            story.append(Paragraph(f"Generated: {date_str}", meta_style))
            story.append(Spacer(1, 30))

        # Create additional styles for different heading levels
        h2_style = ParagraphStyle(
            "CustomH2",
            parent=styles["Heading2"],
            fontSize=self.heading_font_size,
            spaceBefore=15,
            spaceAfter=8,
            fontName=f"{self.font_name}-Bold" if self.font_name == "Helvetica" else self.font_name,
        )

        h3_style = ParagraphStyle(
            "CustomH3",
            parent=styles["Heading3"],
            fontSize=self.font_size + 1,
            spaceBefore=12,
            spaceAfter=6,
            fontName=f"{self.font_name}-Bold" if self.font_name == "Helvetica" else self.font_name,
        )

        # Process content line by line for better markdown handling
        lines = content.split("\n")
        i = 0
        while i < len(lines):
            line = lines[i].rstrip()

            # Skip empty lines
            if not line:
                story.append(Spacer(1, 6))
                i += 1
                continue

            # Handle headings by level
            if line.startswith("### "):
                heading_text = self._convert_markdown_inline(line[4:].strip())
                story.append(Paragraph(heading_text, h3_style))
            elif line.startswith("## "):
                heading_text = self._convert_markdown_inline(line[3:].strip())
                story.append(Paragraph(heading_text, h2_style))
            elif line.startswith("# "):
                heading_text = self._convert_markdown_inline(line[2:].strip())
                story.append(Paragraph(heading_text, heading_style))
            elif line.endswith(":") and len(line) < 100 and not line.startswith(" "):
                # Section heading (text ending with colon)
                story.append(Paragraph(self._escape_xml(line), heading_style))
            elif line.startswith("- ") or line.startswith("* "):
                # Bullet list item
                bullet_text = self._convert_markdown_inline(line[2:].strip())
                story.append(Paragraph(f"\u2022 {bullet_text}", body_style))
            elif line.strip().startswith(tuple(f"{n}. " for n in range(1, 100))):
                # Numbered list item
                parts = line.strip().split(". ", 1)
                if len(parts) == 2:
                    num, text = parts
                    text = self._convert_markdown_inline(text)
                    story.append(Paragraph(f"{num}. {text}", body_style))
            elif line.startswith("```"):
                # Code block - collect all lines until closing ```
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    code_lines.append(lines[i])
                    i += 1
                code_text = "\n".join(code_lines)
                code_style = ParagraphStyle(
                    "Code",
                    parent=styles["Code"],
                    fontSize=9,
                    fontName="Courier",
                    backColor="#f5f5f5",
                    leftIndent=20,
                    rightIndent=20,
                    spaceBefore=10,
                    spaceAfter=10,
                )
                safe_code = self._escape_xml(code_text)
                story.append(Paragraph(f"<pre>{safe_code}</pre>", code_style))
            else:
                # Regular paragraph - convert markdown inline formatting
                safe_para = self._convert_markdown_inline(line)
                story.append(Paragraph(safe_para, body_style))

            i += 1

        # Build PDF
        doc.build(story)

        # Return file size
        return os.path.getsize(output_path)

    def _escape_xml(self, text: str) -> str:
        """Escape XML special characters for ReportLab."""
        return (
            text.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
        )

    def _convert_markdown_inline(self, text: str) -> str:
        """
        Convert markdown inline formatting to ReportLab XML.

        Handles:
        - **bold** -> <b>bold</b>
        - *italic* -> <i>italic</i>
        - `code` -> <font name="Courier">code</font>
        """
        import re

        # First escape XML special characters (but not our formatting)
        text = self._escape_xml(text)

        # Convert **bold** to <b>bold</b>
        text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)

        # Convert *italic* to <i>italic</i> (but not if it's a bullet)
        text = re.sub(r'(?<!\*)\*([^*]+?)\*(?!\*)', r'<i>\1</i>', text)

        # Convert `code` to <font name="Courier">code</font>
        text = re.sub(r'`([^`]+?)`', r'<font name="Courier" size="9">\1</font>', text)

        return text

    def _generate_simple_pdf(self, content: str, output_path: str) -> int:
        """
        Generate a simple text file as fallback when ReportLab is not available.

        Args:
            content: Text content
            output_path: Output file path

        Returns:
            File size in bytes
        """
        # Write as text file with .pdf extension (not a real PDF)
        # This is a fallback for when ReportLab is not installed
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("=" * 60 + "\n")
            f.write("OpenBench Report (Text Format)\n")
            f.write("ReportLab not installed - install with: pip install reportlab\n")
            f.write("=" * 60 + "\n\n")
            f.write(content)

        return os.path.getsize(output_path)


class MarkdownGenerator(OutputGenerator):
    """
    Generate Markdown files.

    Implements the OutputGenerator interface for Markdown output.

    Example:
        ```python
        generator = MarkdownGenerator()
        result = generator.generate(
            content="# My Report\\n\\nThis is content...",
            output_path="report.md"
        )
        ```
    """

    def __init__(self, add_toc: bool = False):
        """
        Initialize Markdown generator.

        Args:
            add_toc: Whether to add table of contents
        """
        self.add_toc = add_toc
        logger.debug(f"MarkdownGenerator initialized (add_toc: {add_toc})")

    @property
    def output_format(self) -> str:
        """Output format identifier."""
        return "markdown"

    def validate(self, content: Any) -> bool:
        """Validate content can be rendered as Markdown."""
        if content is None:
            return False
        return isinstance(content, (str, dict, list)) or hasattr(content, "__str__")

    def _extract_content(self, content: Any) -> str:
        """Extract text content from various input formats."""
        if isinstance(content, str):
            return content

        if isinstance(content, list):
            return "\n".join(f"- {item}" for item in content)

        if hasattr(content, "content"):
            return str(content.content)

        if not isinstance(content, dict):
            return str(content)

        # Handle dict inputs from various layers
        if "content" in content:
            return str(content["content"])

        if "intelligence_output" in content:
            output = content["intelligence_output"]
            if isinstance(output, dict) and "content" in output:
                return str(output["content"])
            return str(output)

        # Convert remaining dict keys to markdown sections
        parts = [
            f"## {key}\n\n{value}"
            for key, value in content.items()
            if key not in ("metadata", "tokens_used", "model")
        ]
        return "\n\n".join(parts) if parts else str(content)

    def generate(
        self,
        content: Any,
        template: Optional[str] = None,
        output_path: str = "output.md",
        title: Optional[str] = None,
        **options,
    ) -> GeneratedOutput:
        """
        Generate Markdown file.

        Args:
            content: Content to render
            template: Template (unused for markdown)
            output_path: Output file path
            title: Document title
            **options: Additional options

        Returns:
            GeneratedOutput with file path and metadata
        """
        logger.info(f"Generating Markdown: {output_path}")

        text_content = self._extract_content(content)

        # Ensure output directory exists
        output_dir = Path(output_path).parent
        if output_dir and str(output_dir) != ".":
            output_dir.mkdir(parents=True, exist_ok=True)

        # Build markdown content
        md_content = ""

        if title:
            md_content += f"# {title}\n\n"
            md_content += f"*Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}*\n\n"
            md_content += "---\n\n"

        md_content += text_content

        # Write file
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(md_content)

        size_bytes = os.path.getsize(output_path)
        logger.info(f"Markdown generated: {output_path} ({size_bytes} bytes)")

        return GeneratedOutput(
            file_path=output_path,
            format=self.output_format,
            size_bytes=size_bytes,
            metadata={
                "title": title,
                "content_length": len(text_content),
                **options,
            },
        )


class PowerPointGenerator(OutputGenerator):
    """
    Generate PowerPoint presentations.

    Implements the OutputGenerator interface for PPTX output.

    Example:
        >>> generator = PowerPointGenerator(template="corporate")
        >>> slides = [{"title": "Intro", "content": "..."}]
        >>> result = generator.generate(content=slides)
    """

    def __init__(self, template: str = "corporate"):
        """
        Initialize PowerPoint generator.

        Args:
            template: Template name for slide design
        """
        self.template = template
        logger.debug(f"PowerPointGenerator initialized (template: {template})")

    @property
    def output_format(self) -> str:
        """Output format identifier."""
        return "pptx"

    def validate(self, content: Any) -> bool:
        """
        Validate that content can be rendered as PowerPoint.

        Args:
            content: Content to validate (should be list of slide dicts)

        Returns:
            True if content is valid slide data
        """
        if content is None:
            return False
        # Accept list of dicts (slides) or dict with slides key
        if isinstance(content, list):
            return True
        if isinstance(content, dict):
            return "slides" in content or len(content) > 0
        return False

    def generate(
        self,
        content: Any,
        template: Optional[str] = None,
        output_path: str = "presentation.pptx",
        **options,
    ) -> GeneratedOutput:
        """
        Generate PowerPoint presentation.

        Args:
            content: Slide content (list of dicts or dict with slides key)
            template: Template override
            output_path: Output file path
            **options: Additional PPTX-specific options

        Returns:
            GeneratedOutput with file path and metadata
        """
        used_template = template or self.template

        # Normalize slides data
        if isinstance(content, list):
            slides = content
        elif isinstance(content, dict) and "slides" in content:
            slides = content["slides"]
        else:
            slides = [content] if content else []

        logger.info(f"Generating {len(slides)} slides: {output_path}")

        # Ensure output directory exists
        output_dir = Path(output_path).parent
        if output_dir and str(output_dir) != ".":
            output_dir.mkdir(parents=True, exist_ok=True)

        # Try to use python-pptx
        try:
            size_bytes = self._generate_with_pptx(slides, output_path, used_template, **options)
        except ImportError:
            logger.warning("python-pptx not installed, creating placeholder file")
            size_bytes = self._generate_placeholder(slides, output_path)

        return GeneratedOutput(
            file_path=output_path,
            format=self.output_format,
            size_bytes=size_bytes,
            metadata={
                "template": used_template,
                "slide_count": len(slides),
                **options,
            },
        )

    def _generate_with_pptx(
        self, slides: List[Any], output_path: str, template: str, **options
    ) -> int:
        """Generate PPTX using python-pptx."""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()

        for slide_data in slides:
            # Use blank layout
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Add title if present
            if isinstance(slide_data, dict):
                title_text = slide_data.get("title", "")
                content_text = slide_data.get("content", str(slide_data))
            else:
                title_text = ""
                content_text = str(slide_data)

            if title_text:
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = title_text
                title_frame.paragraphs[0].font.size = Pt(32)
                title_frame.paragraphs[0].font.bold = True

            # Add content
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(9), Inches(5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.text = content_text

        prs.save(output_path)
        return os.path.getsize(output_path)

    def _generate_placeholder(self, slides: List[Any], output_path: str) -> int:
        """Generate placeholder text file when python-pptx not installed."""
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("PowerPoint Placeholder\n")
            f.write("Install python-pptx: pip install python-pptx\n\n")
            for i, slide in enumerate(slides, 1):
                f.write(f"--- Slide {i} ---\n{slide}\n\n")
        return os.path.getsize(output_path)


class DashboardGenerator(OutputGenerator):
    """
    Generate interactive dashboards.

    Implements the OutputGenerator interface for dashboard output.

    Example:
        >>> generator = DashboardGenerator(framework="streamlit")
        >>> result = generator.generate(content=data, port=8501)
    """

    def __init__(self, framework: str = "streamlit"):
        """
        Initialize dashboard generator.

        Args:
            framework: Dashboard framework ('streamlit', 'dash', 'gradio')
        """
        self.framework = framework
        logger.debug(f"DashboardGenerator initialized (framework: {framework})")

    @property
    def output_format(self) -> str:
        """Output format identifier."""
        return "dashboard"

    def validate(self, content: Any) -> bool:
        """
        Validate that content can be rendered as dashboard.

        Args:
            content: Content to validate

        Returns:
            True if content is valid dashboard data
        """
        if content is None:
            return False
        # Accept dicts, lists, or dataframe-like objects
        return isinstance(content, (dict, list)) or hasattr(content, "to_dict")

    def generate(
        self,
        content: Any,
        template: Optional[str] = None,
        port: int = 8501,
        output_path: Optional[str] = None,
        **options,
    ) -> GeneratedOutput:
        """
        Generate interactive dashboard.

        Args:
            content: Data to visualize in dashboard
            template: Dashboard template/layout
            port: Port to serve dashboard on
            output_path: Path for generated dashboard files
            **options: Additional dashboard-specific options

        Returns:
            GeneratedOutput with dashboard URL and metadata
        """
        dashboard_url = f"http://localhost:{port}"
        file_path = output_path or f"dashboard_{port}"

        logger.info(f"Creating {self.framework} dashboard on port {port}")

        return GeneratedOutput(
            file_path=file_path,
            format=self.output_format,
            size_bytes=0,  # Dashboards are served, not static files
            metadata={
                "framework": self.framework,
                "url": dashboard_url,
                "port": port,
                "template": template,
                **options,
            },
        )


class AudioGenerator(OutputGenerator):
    """
    Generate audio content from text.

    Implements the OutputGenerator interface for audio output.

    Example:
        >>> generator = AudioGenerator(provider="elevenlabs", voice="professional_male")
        >>> result = generator.generate(content="Hello world", output_path="greeting.mp3")
    """

    def __init__(self, provider: str = "elevenlabs", voice: str = "professional_male"):
        """
        Initialize audio generator.

        Args:
            provider: TTS provider ('elevenlabs', 'openai', 'google')
            voice: Voice ID or name to use
        """
        self.provider = provider
        self.voice = voice
        logger.debug(f"AudioGenerator initialized (provider: {provider}, voice: {voice})")

    @property
    def output_format(self) -> str:
        """Output format identifier."""
        return "audio"

    def validate(self, content: Any) -> bool:
        """
        Validate that content can be rendered as audio.

        Args:
            content: Content to validate (should be text or SSML)

        Returns:
            True if content is valid for TTS
        """
        if content is None:
            return False
        # Accept strings or objects with text property
        if isinstance(content, str):
            return len(content.strip()) > 0
        if hasattr(content, "text"):
            return len(str(content.text).strip()) > 0
        return False

    def generate(
        self,
        content: Any,
        template: Optional[str] = None,
        output_path: str = "audio.mp3",
        **options,
    ) -> GeneratedOutput:
        """
        Generate audio from text.

        Args:
            content: Text to convert to speech
            template: Audio style/template (unused, for interface compatibility)
            output_path: Output file path
            **options: Additional audio-specific options (speed, pitch, etc.)

        Returns:
            GeneratedOutput with audio file path and metadata
        """
        # Extract text from content
        if isinstance(content, str):
            text = content
        elif hasattr(content, "text"):
            text = str(content.text)
        else:
            text = str(content)

        logger.info(f"Generating audio: {output_path}")

        # Placeholder for actual TTS generation
        # Estimate size based on text length (rough MP3 estimate)
        estimated_duration_secs = len(text.split()) / 2.5  # ~150 words per minute
        estimated_size = int(estimated_duration_secs * 16000)  # ~128kbps MP3

        return GeneratedOutput(
            file_path=output_path,
            format=self.output_format,
            size_bytes=estimated_size,
            metadata={
                "provider": self.provider,
                "voice": self.voice,
                "text_length": len(text),
                "estimated_duration_secs": estimated_duration_secs,
                **options,
            },
        )
