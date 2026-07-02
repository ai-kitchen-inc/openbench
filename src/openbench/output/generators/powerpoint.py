"""PowerPoint (PPTX) presentation generation using python-pptx."""

from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Any

from openbench.core.abstractions import GeneratedOutput, OutputGenerator

logger = logging.getLogger(__name__)


class PowerPointGenerator(OutputGenerator):
    """
    Generate PowerPoint presentations.

    Implements the OutputGenerator interface for PPTX output.

    Example:
        >>> generator = PowerPointGenerator(template="corporate")
        >>> slides = [{"title": "Intro", "content": "..."}]
        >>> result = generator.generate(content=slides)
    """

    def __init__(self, template: str = "corporate"):
        """
        Initialize PowerPoint generator.

        Args:
            template: Template name for slide design
        """
        self.template = template
        logger.debug(f"PowerPointGenerator initialized (template: {template})")

    @property
    def output_format(self) -> str:
        """Output format identifier."""
        return "pptx"

    def validate(self, content: Any) -> bool:
        """
        Validate that content can be rendered as PowerPoint.

        Args:
            content: Content to validate (should be list of slide dicts)

        Returns:
            True if content is valid slide data
        """
        if content is None:
            return False
        # Accept list of dicts (slides) or dict with slides key
        if isinstance(content, list):
            return True
        if isinstance(content, dict):
            return "slides" in content or len(content) > 0
        return False

    def generate(
        self,
        content: Any,
        template: str | None = None,
        output_path: str = "presentation.pptx",
        **options,
    ) -> GeneratedOutput:
        """
        Generate PowerPoint presentation.

        Args:
            content: Slide content (list of dicts or dict with slides key)
            template: Template override
            output_path: Output file path
            **options: Additional PPTX-specific options

        Returns:
            GeneratedOutput with file path and metadata
        """
        used_template = template or self.template

        # Normalize slides data
        if isinstance(content, list):
            slides = content
        elif isinstance(content, dict) and "slides" in content:
            slides = content["slides"]
        else:
            slides = [content] if content else []

        logger.info(f"Generating {len(slides)} slides: {output_path}")

        # Ensure output directory exists
        output_dir = Path(output_path).parent
        if output_dir and str(output_dir) != ".":
            output_dir.mkdir(parents=True, exist_ok=True)

        # Try to use python-pptx
        try:
            size_bytes = self._generate_with_pptx(slides, output_path, used_template, **options)
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint generation. "
                "Install with: pip install openbench[output]"
            ) from None

        return GeneratedOutput(
            file_path=output_path,
            format=self.output_format,
            size_bytes=size_bytes,
            metadata={
                "template": used_template,
                "slide_count": len(slides),
                **options,
            },
        )

    def _generate_with_pptx(
        self, slides: list[Any], output_path: str, template: str, **options
    ) -> int:
        """Generate PPTX using python-pptx."""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()

        for slide_data in slides:
            # Use blank layout
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Add title if present
            if isinstance(slide_data, dict):
                title_text = slide_data.get("title", "")
                content_text = slide_data.get("content", str(slide_data))
            else:
                title_text = ""
                content_text = str(slide_data)

            if title_text:
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = title_text
                title_frame.paragraphs[0].font.size = Pt(32)
                title_frame.paragraphs[0].font.bold = True

            # Add content
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.text = content_text

        prs.save(output_path)
        return os.path.getsize(output_path)

    def _generate_placeholder(self, slides: list[Any], output_path: str) -> int:
        """Removed: was writing text files with .pptx extension."""
        raise ImportError(
            "python-pptx is required for PowerPoint generation. "
            "Install with: pip install openbench[output]"
        )
